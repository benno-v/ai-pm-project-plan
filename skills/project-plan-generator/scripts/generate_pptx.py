#!/usr/bin/env python3
"""
Generate a client-ready PowerPoint deck from a ``plan`` JSON file.

Design language
---------------
Premium editorial. Muted off-white page, generous white space, large Georgia
display type, Calibri body, a single accent colour used sparingly on progress
fills, today markers, and key numbers. Status tones (ok / warn / risk) are
desaturated so the deck does not read like a traffic light.

Slides produced, in order
-------------------------
  1. Cover                  — project name, client, PM, status-as-of, logo
  2. Executive snapshot     — 4 headline stat cards + status summary
  3. Timeline overview      — phase-level bars on a month-scale rail
  4. Detail Gantt           — task-level rows with progress fills + today line
  5. Milestones             — visual milestone tracker (next six, chronological)
  6. Next 14 days           — what lands in the upcoming fortnight
  7. Review needed          — conditional, only if ``flags`` is non-empty
  8. Handover               — closing card with contact and hand-over line

Every slide carries a muted footer band with an optional client logo, a
project signature line, and a page counter.

Usage
-----
    python generate_pptx.py --plan plan.json --out plan.pptx

Dependencies
------------
    python-pptx>=0.6.21
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from brand import derive_palette, resolve_logo, Palette  # noqa: E402


# ---------------------------------------------------------------------------
# Layout constants (16:9 deck — 13.33" × 7.5")
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PAGE_MARGIN_L = Inches(0.6)
PAGE_MARGIN_R = Inches(0.6)
PAGE_MARGIN_T = Inches(0.55)
PAGE_MARGIN_B = Inches(0.55)

CONTENT_W = SLIDE_W - PAGE_MARGIN_L - PAGE_MARGIN_R

FOOTER_H = Inches(0.35)
FOOTER_TOP = SLIDE_H - FOOTER_H - Inches(0.18)

STATUS_COLOUR_MAP = {
    "achieved": "ok",
    "done": "ok",
    "on_track": "ok",
    "in_progress": "ok",
    "at_risk": "warn",
    "slipping": "warn",
    "delayed": "risk",
    "blocked": "risk",
    "cancelled": "ink_faint",
    "not_started": "ink_faint",
}

STATUS_LABEL_MAP = {
    "achieved": "Achieved",
    "done": "Done",
    "on_track": "On track",
    "in_progress": "In progress",
    "at_risk": "At risk",
    "slipping": "Slipping",
    "delayed": "Delayed",
    "blocked": "Blocked",
    "cancelled": "Cancelled",
    "not_started": "Not started",
}


# ---------------------------------------------------------------------------
# Primitive helpers
# ---------------------------------------------------------------------------

def _parse_date(value: Any) -> Optional[date]:
    if value is None or value == "":
        return None
    if isinstance(value, date) and not isinstance(value, datetime):
        return value
    if isinstance(value, datetime):
        return value.date()
    s = str(value).strip()
    for fmt in ("%Y-%m-%d", "%d %B %Y", "%d %b %Y", "%d/%m/%Y"):
        try:
            return datetime.strptime(s, fmt).date()
        except ValueError:
            continue
    return None


def _hex_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _format_long_date(d: date) -> str:
    return d.strftime("%-d %B %Y") if sys.platform != "win32" else d.strftime("%#d %B %Y")


def _format_short_date(d: date) -> str:
    return d.strftime("%-d %b") if sys.platform != "win32" else d.strftime("%#d %b")


def _week_mondays(start: date, end: date) -> list[date]:
    """Return every Monday on or before ``start`` through on or after ``end``."""
    first_mon = start - timedelta(days=start.weekday())
    last_mon = end - timedelta(days=end.weekday())
    if last_mon < end:
        last_mon = last_mon + timedelta(days=7)
    out: list[date] = []
    cur = first_mon
    while cur <= last_mon:
        out.append(cur)
        cur = cur + timedelta(days=7)
    return out


def _month_starts(start: date, end: date) -> list[date]:
    """Return first-of-month dates covering the span."""
    out = []
    cur = date(start.year, start.month, 1)
    while cur <= end:
        out.append(cur)
        if cur.month == 12:
            cur = date(cur.year + 1, 1, 1)
        else:
            cur = date(cur.year, cur.month + 1, 1)
    return out


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def _set_fill(shape, hex_colour: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_rgb(hex_colour)


def _set_no_fill(shape) -> None:
    shape.fill.background()


def _set_line(shape, hex_colour: Optional[str], width_pt: float = 0.75) -> None:
    if hex_colour is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = _hex_rgb(hex_colour)
    shape.line.width = Pt(width_pt)


def _rect(slide, x, y, w, h, fill_hex: Optional[str] = None,
          line_hex: Optional[str] = None, line_pt: float = 0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill_hex is None:
        _set_no_fill(shp)
    else:
        _set_fill(shp, fill_hex)
    _set_line(shp, line_hex, line_pt)
    shp.shadow.inherit = False
    return shp


def _round_rect(slide, x, y, w, h, fill_hex: Optional[str] = None,
                line_hex: Optional[str] = None, corner: float = 0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill_hex is None:
        _set_no_fill(shp)
    else:
        _set_fill(shp, fill_hex)
    _set_line(shp, line_hex)
    shp.shadow.inherit = False
    # Set corner radius via the "adj" handle — value 0..0.5 as a fraction of min(w,h)/2
    try:
        shp.adjustments[0] = corner
    except (IndexError, AttributeError):
        pass
    return shp


def _ellipse(slide, x, y, w, h, fill_hex: Optional[str] = None,
             line_hex: Optional[str] = None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill_hex is None:
        _set_no_fill(shp)
    else:
        _set_fill(shp, fill_hex)
    _set_line(shp, line_hex)
    shp.shadow.inherit = False
    return shp


def _line(slide, x1, y1, x2, y2, colour_hex: str, width_pt: float = 1.0,
          dashed: bool = False):
    shp = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    shp.line.color.rgb = _hex_rgb(colour_hex)
    shp.line.width = Pt(width_pt)
    if dashed:
        ln = shp.line._get_or_add_ln()
        prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        # Remove any existing prstDash first
        for existing in ln.findall(qn("a:prstDash")):
            ln.remove(existing)
        ln.append(prstDash)
    shp.shadow.inherit = False
    return shp


def _textbox(slide, x, y, w, h, text: str, *,
             font_name: str, size_pt: float,
             colour_hex: str = "#1A1A1A",
             bold: bool = False, italic: bool = False,
             align: str = "left",
             anchor: str = "top",
             line_spacing: float = 1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    anchor_map = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_rgb(colour_hex)
    return tb


def _add_logo(slide, logo_path: Path, x, y, h):
    """Add a logo keeping aspect ratio, constrained to ``h`` tall."""
    try:
        pic = slide.shapes.add_picture(str(logo_path), x, y, height=h)
    except Exception:
        return None
    return pic


# ---------------------------------------------------------------------------
# Footer — applied to every slide
# ---------------------------------------------------------------------------

def _footer(slide, *, palette: Palette, project: dict, logo_path: Optional[Path],
            page_number: int, page_total: int):
    # Hairline rule above the footer
    rule_y = FOOTER_TOP
    _line(slide, PAGE_MARGIN_L, rule_y,
          SLIDE_W - PAGE_MARGIN_R, rule_y, palette.hairline, width_pt=0.5)

    # Logo on the left
    band_top = rule_y + Inches(0.08)
    band_h = FOOTER_H - Inches(0.08)
    logo_right_edge = PAGE_MARGIN_L
    if logo_path is not None:
        pic = _add_logo(slide, logo_path, PAGE_MARGIN_L, band_top, band_h)
        if pic is not None:
            logo_right_edge = PAGE_MARGIN_L + pic.width + Inches(0.18)

    # Centre: project signature line
    signature = " · ".join(
        part for part in [
            project.get("client"),
            project.get("name"),
            "Status as of " + _format_long_date(_parse_date(project.get("status_as_of")))
            if _parse_date(project.get("status_as_of")) else None,
        ] if part
    )
    _textbox(
        slide,
        logo_right_edge, band_top,
        SLIDE_W - logo_right_edge - PAGE_MARGIN_R - Inches(0.8),
        band_h,
        signature,
        font_name=palette.body_font, size_pt=9,
        colour_hex=palette.ink_faint, italic=True,
        align="left", anchor="middle",
    )

    # Page number on the right
    _textbox(
        slide,
        SLIDE_W - PAGE_MARGIN_R - Inches(0.8), band_top,
        Inches(0.8), band_h,
        f"{page_number:02d} / {page_total:02d}",
        font_name=palette.body_font, size_pt=9,
        colour_hex=palette.ink_faint,
        align="right", anchor="middle",
    )


def _page_background(slide, palette: Palette) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_rgb(palette.page)


# ---------------------------------------------------------------------------
# Small composed widgets
# ---------------------------------------------------------------------------

def _eyebrow(slide, x, y, w, text: str, palette: Palette):
    _textbox(
        slide, x, y, w, Inches(0.22),
        text.upper(),
        font_name=palette.body_font, size_pt=10,
        colour_hex=palette.accent, bold=True,
        align="left",
    )


def _section_title(slide, x, y, w, text: str, palette: Palette,
                   size_pt: float = 26):
    _textbox(
        slide, x, y, w, Inches(0.6),
        text,
        font_name=palette.display_font, size_pt=size_pt,
        colour_hex=palette.primary, bold=True,
        align="left",
    )


def _kicker_rule(slide, x, y, palette: Palette, width: Inches = Inches(0.6)):
    _rect(slide, x, y, width, Emu(25400), fill_hex=palette.accent,
          line_hex=None)


def _status_chip(slide, x, y, w, h, status: str, palette: Palette):
    colour_key = STATUS_COLOUR_MAP.get(status, "ink_faint")
    bg_key = colour_key + "_soft" if hasattr(palette, colour_key + "_soft") else "hairline"
    fg_hex = getattr(palette, colour_key, palette.ink_faint)
    bg_hex = getattr(palette, bg_key, palette.hairline)
    _round_rect(slide, x, y, w, h, fill_hex=bg_hex, line_hex=None, corner=0.45)
    _textbox(
        slide, x, y, w, h,
        STATUS_LABEL_MAP.get(status, status.replace("_", " ").title()),
        font_name=palette.body_font, size_pt=9,
        colour_hex=fg_hex, bold=True,
        align="center", anchor="middle",
    )


def _stat_card(slide, x, y, w, h, *, label: str, value: str, sub: str,
               palette: Palette, accent: bool = False):
    _round_rect(slide, x, y, w, h,
                fill_hex=palette.card, line_hex=palette.hairline,
                corner=0.05)

    # Label (eyebrow)
    _textbox(
        slide, x + Inches(0.25), y + Inches(0.22),
        w - Inches(0.5), Inches(0.22),
        label.upper(),
        font_name=palette.body_font, size_pt=9,
        colour_hex=palette.ink_faint, bold=True,
    )

    # Value (big display)
    _textbox(
        slide, x + Inches(0.25), y + Inches(0.48),
        w - Inches(0.5), Inches(0.75),
        value,
        font_name=palette.display_font, size_pt=34,
        colour_hex=palette.accent if accent else palette.primary,
        bold=True, align="left",
    )

    # Sub line
    _textbox(
        slide, x + Inches(0.25), y + h - Inches(0.45),
        w - Inches(0.5), Inches(0.25),
        sub,
        font_name=palette.body_font, size_pt=10,
        colour_hex=palette.ink_soft, italic=True,
    )


# ---------------------------------------------------------------------------
# Summary computations
# ---------------------------------------------------------------------------

def _compute_summary(plan: dict) -> dict:
    tasks = plan.get("tasks", [])
    milestones = plan.get("milestones", [])

    # Overall progress — duration-weighted average of progress
    total_days = 0.0
    weighted = 0.0
    for t in tasks:
        d = float(t.get("duration_days") or 0)
        p = float(t.get("progress") or 0)
        total_days += d
        weighted += d * p
    overall_progress = (weighted / total_days) if total_days else 0.0

    at_risk_tasks = sum(1 for t in tasks if t.get("status") in ("at_risk", "slipping", "delayed", "blocked"))

    achieved = sum(1 for m in milestones if m.get("status") == "achieved")
    total_ms = len(milestones)
    next_milestone = None
    status_d = _parse_date(plan.get("project", {}).get("status_as_of")) or date.today()
    for m in sorted(milestones, key=lambda x: _parse_date(x.get("date")) or date.max):
        md = _parse_date(m.get("date"))
        if md and md >= status_d and m.get("status") != "achieved":
            next_milestone = m
            break
    if next_milestone is None and milestones:
        for m in sorted(milestones, key=lambda x: _parse_date(x.get("date")) or date.max):
            if m.get("status") != "achieved":
                next_milestone = m
                break

    proj = plan.get("project", {})
    start = _parse_date(proj.get("start"))
    end = _parse_date(proj.get("end"))

    schedule_health = "On track"
    if plan.get("flags"):
        schedule_health = "Watch"
    if at_risk_tasks >= 3:
        schedule_health = "Action"

    return {
        "overall_progress": overall_progress,
        "achieved": achieved,
        "total_ms": total_ms,
        "next_milestone": next_milestone,
        "schedule_health": schedule_health,
        "at_risk_tasks": at_risk_tasks,
        "start": start,
        "end": end,
        "status_as_of": status_d,
    }


# ---------------------------------------------------------------------------
# Slide: Cover
# ---------------------------------------------------------------------------

def build_cover(prs: Presentation, plan: dict, palette: Palette,
                logo_path: Optional[Path], page_number: int, page_total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)

    proj = plan.get("project", {})

    # Accent kicker bar top-left
    _rect(slide, PAGE_MARGIN_L, Inches(0.55), Inches(0.6), Inches(0.05),
          fill_hex=palette.accent, line_hex=None)

    # Eyebrow
    _textbox(
        slide, PAGE_MARGIN_L, Inches(0.75),
        CONTENT_W, Inches(0.3),
        (proj.get("client", "") + " · Project plan").upper(),
        font_name=palette.body_font, size_pt=11,
        colour_hex=palette.ink_faint, bold=True,
    )

    # Title — very large
    _textbox(
        slide, PAGE_MARGIN_L, Inches(1.4),
        CONTENT_W, Inches(2.4),
        proj.get("name", "Project plan"),
        font_name=palette.display_font, size_pt=54,
        colour_hex=palette.primary, bold=True,
        line_spacing=1.05,
    )

    # Subtitle — italic status summary, truncated
    summary = plan.get("status_summary") or ""
    if len(summary) > 220:
        summary = summary[:217] + "…"
    _textbox(
        slide, PAGE_MARGIN_L, Inches(3.9),
        CONTENT_W * 0.85, Inches(1.2),
        summary,
        font_name=palette.display_font, size_pt=18,
        colour_hex=palette.ink_soft, italic=True,
        line_spacing=1.4,
    )

    # Detail rows at the bottom left
    rows_top = Inches(5.6)
    col_w = Inches(3.2)
    details = [
        ("CLIENT", proj.get("client", "—")),
        ("PROJECT MANAGER", proj.get("pm", "—")),
        ("STATUS AS OF",
         _format_long_date(_parse_date(proj.get("status_as_of")))
         if _parse_date(proj.get("status_as_of")) else "—"),
    ]
    for i, (label, value) in enumerate(details):
        x = PAGE_MARGIN_L + col_w * i
        _textbox(
            slide, x, rows_top, col_w - Inches(0.2), Inches(0.22),
            label,
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.ink_faint, bold=True,
        )
        _textbox(
            slide, x, rows_top + Inches(0.28), col_w - Inches(0.2), Inches(0.35),
            value,
            font_name=palette.display_font, size_pt=15,
            colour_hex=palette.primary,
        )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Executive snapshot
# ---------------------------------------------------------------------------

def build_exec_snapshot(prs: Presentation, plan: dict, palette: Palette,
                        logo_path: Optional[Path], page_number: int,
                        page_total: int, summary: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Executive snapshot", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "Where the project stands today", palette)

    # Four stat cards
    cards_top = Inches(1.8)
    card_h = Inches(1.9)
    gap = Inches(0.22)
    card_w = (CONTENT_W - gap * 3) / 4

    next_ms = summary["next_milestone"]
    next_ms_label = "—"
    next_ms_sub = "No upcoming milestone"
    if next_ms:
        md = _parse_date(next_ms.get("date"))
        next_ms_label = _format_short_date(md) if md else "—"
        next_ms_sub = next_ms.get("name", "")

    stats = [
        {
            "label": "Overall complete",
            "value": f"{int(round(summary['overall_progress'] * 100))}%",
            "sub": "Duration-weighted average",
            "accent": True,
        },
        {
            "label": "Milestones",
            "value": f"{summary['achieved']} / {summary['total_ms']}",
            "sub": "Achieved to date",
            "accent": False,
        },
        {
            "label": "Next milestone",
            "value": next_ms_label,
            "sub": next_ms_sub,
            "accent": False,
        },
        {
            "label": "Schedule health",
            "value": summary["schedule_health"],
            "sub": f"{summary['at_risk_tasks']} task(s) at risk"
                   if summary['at_risk_tasks']
                   else "No at-risk tasks",
            "accent": False,
        },
    ]

    for i, s in enumerate(stats):
        x = PAGE_MARGIN_L + (card_w + gap) * i
        _stat_card(slide, x, cards_top, card_w, card_h,
                   label=s["label"], value=s["value"], sub=s["sub"],
                   palette=palette, accent=s["accent"])

    # Status summary pull-quote
    quote_top = cards_top + card_h + Inches(0.5)
    quote_h = Inches(1.8)
    # Left accent line
    _rect(slide, PAGE_MARGIN_L, quote_top,
          Emu(25400 * 2), quote_h, fill_hex=palette.accent, line_hex=None)

    _textbox(
        slide, PAGE_MARGIN_L + Inches(0.35), quote_top,
        CONTENT_W - Inches(0.7), quote_h,
        plan.get("status_summary", ""),
        font_name=palette.display_font, size_pt=18,
        colour_hex=palette.ink, italic=True,
        line_spacing=1.4, anchor="middle",
    )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Timeline overview (phase-level)
# ---------------------------------------------------------------------------

def build_timeline_overview(prs: Presentation, plan: dict, palette: Palette,
                            logo_path: Optional[Path], page_number: int,
                            page_total: int, summary: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Timeline", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "Phases at a glance", palette)

    phases = plan.get("phases", [])
    if not phases:
        _textbox(
            slide, PAGE_MARGIN_L, Inches(2.2),
            CONTENT_W, Inches(0.5),
            "No phases defined in the plan.",
            font_name=palette.body_font, size_pt=12,
            colour_hex=palette.ink_faint, italic=True,
        )
        _footer(slide, palette=palette, project=proj, logo_path=logo_path,
                page_number=page_number, page_total=page_total)
        return slide

    # Span of the plan
    start = summary["start"] or min(_parse_date(p.get("start")) for p in phases)
    end = summary["end"] or max(_parse_date(p.get("end")) for p in phases)
    status_d = summary["status_as_of"]
    total_days = max(1, (end - start).days)

    # Rail layout
    chart_left = PAGE_MARGIN_L + Inches(1.8)   # label column on the left
    chart_right = SLIDE_W - PAGE_MARGIN_R - Inches(0.3)
    chart_w = chart_right - chart_left

    chart_top = Inches(1.8)
    header_h = Inches(0.45)
    chart_bottom = FOOTER_TOP - Inches(0.4)

    # Month header
    months = _month_starts(start, end)
    # Hairline separator above rail
    _line(slide, chart_left, chart_top + header_h,
          chart_right, chart_top + header_h, palette.hairline, 0.5)

    for i, m in enumerate(months):
        days_from_start = (m - start).days
        x = chart_left + int(chart_w * (days_from_start / total_days))
        _textbox(
            slide, x, chart_top, Inches(1.5), Inches(0.3),
            m.strftime("%b %Y"),
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.ink_faint, bold=True,
        )
        if i > 0:
            _line(slide, x, chart_top + header_h,
                  x, chart_bottom, palette.hairline, 0.4)

    # Phase rows
    row_top = chart_top + header_h + Inches(0.25)
    row_h = Inches(0.55)
    row_gap = Inches(0.16)

    # Clamp number of phase rows so they fit
    available_h = chart_bottom - row_top
    max_rows = max(1, int(available_h / (row_h + row_gap)))
    phases_to_show = phases[:max_rows]

    for i, ph in enumerate(phases_to_show):
        y = row_top + (row_h + row_gap) * i
        ps = _parse_date(ph.get("start"))
        pe = _parse_date(ph.get("end"))
        if not ps or not pe:
            continue
        x1 = chart_left + int(chart_w * ((ps - start).days / total_days))
        x2 = chart_left + int(chart_w * ((pe - start).days / total_days))
        bar_w = max(Inches(0.2), x2 - x1)

        # Phase label on the left
        _textbox(
            slide, PAGE_MARGIN_L, y + Inches(0.08),
            chart_left - PAGE_MARGIN_L - Inches(0.2), Inches(0.3),
            ph.get("name", ""),
            font_name=palette.body_font, size_pt=11,
            colour_hex=palette.ink, bold=True,
        )
        # Phase span chip
        _textbox(
            slide, PAGE_MARGIN_L, y + Inches(0.3),
            chart_left - PAGE_MARGIN_L - Inches(0.2), Inches(0.24),
            f"{_format_short_date(ps)} — {_format_short_date(pe)}",
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.ink_faint, italic=True,
        )

        # Bar background
        _round_rect(slide, x1, y + Inches(0.08),
                    bar_w, Inches(0.38),
                    fill_hex=palette.primary_soft, line_hex=None,
                    corner=0.45)

        # Progress fill — compute from tasks in this phase
        phase_tasks = [t for t in plan.get("tasks", []) if t.get("phase_id") == ph.get("id")]
        if phase_tasks:
            total_d = sum(float(t.get("duration_days") or 0) for t in phase_tasks)
            weighted = sum(float(t.get("duration_days") or 0) * float(t.get("progress") or 0)
                           for t in phase_tasks)
            pct = (weighted / total_d) if total_d else 0.0
        else:
            pct = 0.0
        if pct > 0:
            fill_w = int(bar_w * min(1.0, pct))
            if fill_w > 0:
                _round_rect(slide, x1, y + Inches(0.08),
                            fill_w, Inches(0.38),
                            fill_hex=palette.accent, line_hex=None,
                            corner=0.45)
        # Phase name label on the bar (only if the bar is wide enough)
        if bar_w > Inches(1.2):
            _textbox(
                slide, x1 + Inches(0.12), y + Inches(0.11),
                bar_w - Inches(0.24), Inches(0.3),
                f"{int(round(pct * 100))}% complete",
                font_name=palette.body_font, size_pt=9,
                colour_hex=palette.page if pct > 0.45 else palette.ink,
                bold=True, anchor="middle",
            )

    # Today marker
    if start <= status_d <= end:
        today_x = chart_left + int(chart_w * ((status_d - start).days / total_days))
        _line(slide, today_x, chart_top + header_h,
              today_x, chart_bottom, palette.accent, 1.5, dashed=True)
        _textbox(
            slide, today_x - Inches(0.35), chart_bottom + Inches(0.03),
            Inches(0.8), Inches(0.22),
            "TODAY",
            font_name=palette.body_font, size_pt=8,
            colour_hex=palette.accent, bold=True,
            align="center",
        )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Detail Gantt (task-level)
# ---------------------------------------------------------------------------

def build_detail_gantt(prs: Presentation, plan: dict, palette: Palette,
                       logo_path: Optional[Path], page_number: int,
                       page_total: int, summary: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Detail schedule", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "Task-level Gantt", palette)

    tasks = plan.get("tasks", [])
    phases = plan.get("phases", [])
    if not tasks:
        _textbox(
            slide, PAGE_MARGIN_L, Inches(2.2),
            CONTENT_W, Inches(0.5),
            "No tasks to display.",
            font_name=palette.body_font, size_pt=12,
            colour_hex=palette.ink_faint, italic=True,
        )
        _footer(slide, palette=palette, project=proj, logo_path=logo_path,
                page_number=page_number, page_total=page_total)
        return slide

    # Span
    start = summary["start"] or min(filter(None, (_parse_date(t.get("start")) for t in tasks)))
    end = summary["end"] or max(filter(None, (_parse_date(t.get("end")) for t in tasks)))
    status_d = summary["status_as_of"]
    total_days = max(1, (end - start).days)

    # Layout
    chart_left = PAGE_MARGIN_L + Inches(2.7)    # left column: task name + owner
    chart_right = SLIDE_W - PAGE_MARGIN_R - Inches(0.3)
    chart_w = chart_right - chart_left

    chart_top = Inches(1.75)
    header_h = Inches(0.35)
    chart_bottom = FOOTER_TOP - Inches(0.3)

    # Month ticks
    months = _month_starts(start, end)
    _line(slide, chart_left, chart_top + header_h,
          chart_right, chart_top + header_h, palette.hairline, 0.5)
    for i, m in enumerate(months):
        days = (m - start).days
        x = chart_left + int(chart_w * (days / total_days))
        _textbox(
            slide, x, chart_top, Inches(1.2), Inches(0.28),
            m.strftime("%b").upper(),
            font_name=palette.body_font, size_pt=8,
            colour_hex=palette.ink_faint, bold=True,
        )
        if i > 0:
            _line(slide, x, chart_top + header_h,
                  x, chart_bottom, palette.hairline, 0.4)

    # Sort tasks by phase order then start date
    phase_order = {p.get("id"): idx for idx, p in enumerate(phases)}
    tasks_sorted = sorted(
        tasks,
        key=lambda t: (phase_order.get(t.get("phase_id"), 99),
                       _parse_date(t.get("start")) or date.max),
    )

    available_h = chart_bottom - chart_top - header_h
    max_rows = max(5, int(available_h / Inches(0.32)))
    tasks_to_show = tasks_sorted[:max_rows]

    row_h = available_h / max(1, len(tasks_to_show))
    if row_h > Inches(0.45):
        row_h = Inches(0.45)
    bar_h = Inches(0.20)

    current_phase_id = None
    for i, t in enumerate(tasks_to_show):
        y = chart_top + header_h + row_h * i + (row_h - bar_h) / 2

        # Subtle phase band (row shading alternates by phase)
        if t.get("phase_id") != current_phase_id:
            current_phase_id = t.get("phase_id")

        ts = _parse_date(t.get("start"))
        te = _parse_date(t.get("end"))
        if not ts or not te:
            continue

        x1 = chart_left + int(chart_w * ((ts - start).days / total_days))
        x2 = chart_left + int(chart_w * ((te - start).days / total_days))
        bar_w = max(Inches(0.08), x2 - x1)

        # Left: task name
        _textbox(
            slide, PAGE_MARGIN_L, chart_top + header_h + row_h * i,
            chart_left - PAGE_MARGIN_L - Inches(0.2), row_h,
            t.get("name", ""),
            font_name=palette.body_font, size_pt=9.5,
            colour_hex=palette.ink, bold=False,
            anchor="middle",
        )
        # Owner (small caps to the right of the task name)
        # — fit into same row via the right side of the name column
        _textbox(
            slide, PAGE_MARGIN_L + Inches(1.75),
            chart_top + header_h + row_h * i,
            Inches(0.9), row_h,
            t.get("owner_role", ""),
            font_name=palette.body_font, size_pt=8,
            colour_hex=palette.ink_faint, italic=True,
            align="right", anchor="middle",
        )

        # Bar background
        status = t.get("status", "not_started")
        bg_key = {
            "at_risk": "warn_soft",
            "slipping": "warn_soft",
            "delayed": "risk_soft",
            "blocked": "risk_soft",
        }.get(status, "primary_soft")
        _round_rect(slide, x1, y, bar_w, bar_h,
                    fill_hex=getattr(palette, bg_key), line_hex=None,
                    corner=0.45)

        # Progress fill
        pct = float(t.get("progress") or 0)
        fg_key = "warn" if status in ("at_risk", "slipping") else \
                 "risk" if status in ("delayed", "blocked") else "accent"
        if pct > 0:
            fill_w = int(bar_w * min(1.0, pct))
            if fill_w > 0:
                _round_rect(slide, x1, y, fill_w, bar_h,
                            fill_hex=getattr(palette, fg_key), line_hex=None,
                            corner=0.45)

    # Today marker
    if start <= status_d <= end:
        today_x = chart_left + int(chart_w * ((status_d - start).days / total_days))
        _line(slide, today_x, chart_top + header_h,
              today_x, chart_bottom, palette.accent, 1.5, dashed=True)
        _textbox(
            slide, today_x - Inches(0.35), chart_top - Inches(0.06),
            Inches(0.8), Inches(0.22),
            "TODAY",
            font_name=palette.body_font, size_pt=8,
            colour_hex=palette.accent, bold=True,
            align="center",
        )

    # Mini legend bottom-left
    legend_y = chart_bottom + Inches(0.05)
    _textbox(
        slide, PAGE_MARGIN_L, legend_y,
        Inches(4), Inches(0.2),
        "Shaded bar: remaining   ·   Accent fill: complete   ·   Warm fill: at risk",
        font_name=palette.body_font, size_pt=8,
        colour_hex=palette.ink_faint, italic=True,
    )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Milestones
# ---------------------------------------------------------------------------

def build_milestones(prs: Presentation, plan: dict, palette: Palette,
                     logo_path: Optional[Path], page_number: int,
                     page_total: int, summary: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Milestones", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "Milestone tracker", palette)

    milestones = sorted(
        plan.get("milestones", []),
        key=lambda m: _parse_date(m.get("date")) or date.max,
    )
    if not milestones:
        _textbox(
            slide, PAGE_MARGIN_L, Inches(2.2),
            CONTENT_W, Inches(0.5),
            "No milestones defined.",
            font_name=palette.body_font, size_pt=12,
            colour_hex=palette.ink_faint, italic=True,
        )
        _footer(slide, palette=palette, project=proj, logo_path=logo_path,
                page_number=page_number, page_total=page_total)
        return slide

    # Horizontal dotted track with dots per milestone
    track_top = Inches(1.95)
    track_left = PAGE_MARGIN_L + Inches(0.2)
    track_right = SLIDE_W - PAGE_MARGIN_R - Inches(0.2)
    track_w = track_right - track_left

    dates = [_parse_date(m.get("date")) for m in milestones if _parse_date(m.get("date"))]
    if not dates:
        _footer(slide, palette=palette, project=proj, logo_path=logo_path,
                page_number=page_number, page_total=page_total)
        return slide
    t_start = min(dates)
    t_end = max(dates)
    total_days = max(1, (t_end - t_start).days)

    # Line
    _line(slide, track_left, track_top, track_right, track_top,
          palette.hairline, 1.0)

    # Cards grid below the track
    cards_top = track_top + Inches(0.55)
    cards_bottom = FOOTER_TOP - Inches(0.3)
    cards_h = cards_bottom - cards_top

    # Pick up to 6 milestones to feature (prioritise next upcoming + any at-risk)
    status_d = summary["status_as_of"]
    upcoming = [m for m in milestones if (_parse_date(m.get("date")) or date.min) >= status_d]
    achieved = [m for m in milestones if m.get("status") == "achieved"]
    feature_pool: list[dict] = []
    # Always include latest achieved (most recent)
    if achieved:
        feature_pool.append(achieved[-1])
    # Then upcoming
    feature_pool.extend(upcoming)
    # Dedupe preserving order
    seen = set()
    feature = []
    for m in feature_pool:
        mid = m.get("id")
        if mid in seen:
            continue
        seen.add(mid)
        feature.append(m)
        if len(feature) >= 6:
            break
    if not feature:
        feature = milestones[:6]

    # Draw all milestone dots on the track (all, not just feature)
    for m in milestones:
        md = _parse_date(m.get("date"))
        if not md:
            continue
        dot_x = track_left + int(track_w * ((md - t_start).days / total_days)) - Inches(0.07)
        status = m.get("status", "on_track")
        colour_key = STATUS_COLOUR_MAP.get(status, "ink_faint")
        fill_hex = getattr(palette, colour_key, palette.ink_faint)
        _ellipse(slide, dot_x, track_top - Inches(0.07),
                 Inches(0.14), Inches(0.14),
                 fill_hex=fill_hex, line_hex=palette.page)

    # Card grid — up to 3 cols x 2 rows
    cols = 3
    rows = max(1, (len(feature) + cols - 1) // cols)
    gap = Inches(0.25)
    card_w = (CONTENT_W - gap * (cols - 1)) / cols
    card_h_fit = (cards_h - gap * (rows - 1)) / max(1, rows)
    card_h = min(Inches(2.2), card_h_fit)

    for i, m in enumerate(feature):
        r = i // cols
        c = i % cols
        x = PAGE_MARGIN_L + (card_w + gap) * c
        y = cards_top + (card_h + gap) * r
        _round_rect(slide, x, y, card_w, card_h,
                    fill_hex=palette.card, line_hex=palette.hairline,
                    corner=0.05)
        # Date eyebrow
        md = _parse_date(m.get("date"))
        _textbox(
            slide, x + Inches(0.25), y + Inches(0.18),
            card_w - Inches(0.5), Inches(0.24),
            (_format_long_date(md).upper() if md else ""),
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.accent, bold=True,
        )
        # Name
        _textbox(
            slide, x + Inches(0.25), y + Inches(0.45),
            card_w - Inches(0.5), Inches(0.5),
            m.get("name", ""),
            font_name=palette.display_font, size_pt=15,
            colour_hex=palette.primary, bold=True,
            line_spacing=1.2,
        )
        # Description
        desc = m.get("description", "")
        _textbox(
            slide, x + Inches(0.25), y + Inches(1.0),
            card_w - Inches(0.5), card_h - Inches(1.4),
            desc,
            font_name=palette.body_font, size_pt=10,
            colour_hex=palette.ink_soft,
            line_spacing=1.3,
        )
        # Status chip bottom-left
        _status_chip(
            slide, x + Inches(0.25), y + card_h - Inches(0.45),
            Inches(0.95), Inches(0.28),
            m.get("status", "on_track"), palette,
        )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Next 14 days
# ---------------------------------------------------------------------------

def build_next_14_days(prs: Presentation, plan: dict, palette: Palette,
                       logo_path: Optional[Path], page_number: int,
                       page_total: int, summary: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Next fortnight", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "What lands in the next 14 days", palette)

    status_d = summary["status_as_of"]
    horizon = status_d + timedelta(days=14)

    # Tasks finishing in the window
    window_tasks = []
    for t in plan.get("tasks", []):
        te = _parse_date(t.get("end"))
        if te and status_d <= te <= horizon:
            window_tasks.append(t)
    window_tasks.sort(key=lambda t: _parse_date(t.get("end")) or date.max)

    # Milestones in the window
    window_ms = []
    for m in plan.get("milestones", []):
        md = _parse_date(m.get("date"))
        if md and status_d <= md <= horizon:
            window_ms.append(m)
    window_ms.sort(key=lambda m: _parse_date(m.get("date")) or date.max)

    # Two columns
    col_w = (CONTENT_W - Inches(0.4)) / 2
    col_top = Inches(1.85)
    col_h = FOOTER_TOP - col_top - Inches(0.3)

    # Left column — Tasks ending
    _textbox(
        slide, PAGE_MARGIN_L, col_top,
        col_w, Inches(0.35),
        "Tasks completing",
        font_name=palette.display_font, size_pt=16,
        colour_hex=palette.primary, bold=True,
    )
    _line(slide, PAGE_MARGIN_L, col_top + Inches(0.4),
          PAGE_MARGIN_L + col_w, col_top + Inches(0.4),
          palette.hairline, 0.5)

    list_top = col_top + Inches(0.55)
    row_h = Inches(0.5)
    max_rows = max(1, int((col_h - Inches(0.55)) / row_h))
    for i, t in enumerate(window_tasks[:max_rows]):
        y = list_top + row_h * i
        te = _parse_date(t.get("end"))
        # Date pill on the left
        _textbox(
            slide, PAGE_MARGIN_L, y,
            Inches(1.1), row_h,
            _format_short_date(te) if te else "—",
            font_name=palette.display_font, size_pt=14,
            colour_hex=palette.accent, bold=True,
            anchor="middle",
        )
        _textbox(
            slide, PAGE_MARGIN_L + Inches(1.15), y,
            col_w - Inches(1.15), Inches(0.25),
            t.get("name", ""),
            font_name=palette.body_font, size_pt=11,
            colour_hex=palette.ink, bold=True,
            anchor="top",
        )
        _textbox(
            slide, PAGE_MARGIN_L + Inches(1.15), y + Inches(0.24),
            col_w - Inches(1.15), Inches(0.25),
            t.get("owner_role", ""),
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.ink_faint, italic=True,
        )

    if not window_tasks:
        _textbox(
            slide, PAGE_MARGIN_L, list_top,
            col_w, Inches(0.5),
            "No tasks complete in this window.",
            font_name=palette.body_font, size_pt=11,
            colour_hex=palette.ink_faint, italic=True,
        )

    # Right column — Milestones
    right_x = PAGE_MARGIN_L + col_w + Inches(0.4)
    _textbox(
        slide, right_x, col_top,
        col_w, Inches(0.35),
        "Milestones in window",
        font_name=palette.display_font, size_pt=16,
        colour_hex=palette.primary, bold=True,
    )
    _line(slide, right_x, col_top + Inches(0.4),
          right_x + col_w, col_top + Inches(0.4),
          palette.hairline, 0.5)

    for i, m in enumerate(window_ms[:max_rows]):
        y = list_top + row_h * i
        md = _parse_date(m.get("date"))
        _textbox(
            slide, right_x, y,
            Inches(1.1), row_h,
            _format_short_date(md) if md else "—",
            font_name=palette.display_font, size_pt=14,
            colour_hex=palette.accent, bold=True,
            anchor="middle",
        )
        _textbox(
            slide, right_x + Inches(1.15), y,
            col_w - Inches(1.15) - Inches(1.0), Inches(0.25),
            m.get("name", ""),
            font_name=palette.body_font, size_pt=11,
            colour_hex=palette.ink, bold=True,
        )
        _textbox(
            slide, right_x + Inches(1.15), y + Inches(0.24),
            col_w - Inches(1.15) - Inches(1.0), Inches(0.25),
            m.get("owner_role", ""),
            font_name=palette.body_font, size_pt=9,
            colour_hex=palette.ink_faint, italic=True,
        )
        _status_chip(
            slide, right_x + col_w - Inches(1.0), y + Inches(0.1),
            Inches(0.95), Inches(0.26),
            m.get("status", "on_track"), palette,
        )

    if not window_ms:
        _textbox(
            slide, right_x, list_top,
            col_w, Inches(0.5),
            "No milestones in this window.",
            font_name=palette.body_font, size_pt=11,
            colour_hex=palette.ink_faint, italic=True,
        )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Review needed (conditional)
# ---------------------------------------------------------------------------

def build_review_needed(prs: Presentation, plan: dict, palette: Palette,
                        logo_path: Optional[Path], page_number: int,
                        page_total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    _eyebrow(slide, PAGE_MARGIN_L, PAGE_MARGIN_T, CONTENT_W,
             "Review needed", palette)
    _section_title(slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(0.3),
                   CONTENT_W, "Items that need a decision", palette)

    _textbox(
        slide, PAGE_MARGIN_L, PAGE_MARGIN_T + Inches(1.0),
        CONTENT_W, Inches(0.3),
        "Bring these to the project review — each one needs a "
        "client decision or an escalation to keep the plan on track.",
        font_name=palette.body_font, size_pt=11,
        colour_hex=palette.ink_soft, italic=True,
    )

    flags = plan.get("flags", [])
    flags_top = Inches(2.7)
    card_h = Inches(1.25)
    gap = Inches(0.2)
    available_h = FOOTER_TOP - flags_top - Inches(0.2)
    max_flags = max(1, int(available_h / (card_h + gap)))
    flags_to_show = flags[:max_flags]

    for i, f in enumerate(flags_to_show):
        y = flags_top + (card_h + gap) * i
        _round_rect(slide, PAGE_MARGIN_L, y, CONTENT_W, card_h,
                    fill_hex=palette.warn_soft, line_hex=palette.warn,
                    corner=0.04)
        # Left bar
        _rect(slide, PAGE_MARGIN_L, y, Inches(0.08), card_h,
              fill_hex=palette.warn, line_hex=None)

        scope = (f.get("scope") or "").upper()
        ref = f.get("ref", "")
        _textbox(
            slide, PAGE_MARGIN_L + Inches(0.35), y + Inches(0.18),
            Inches(3), Inches(0.3),
            f"{scope} · {ref}" if scope or ref else "",
            font_name=palette.body_font, size_pt=10,
            colour_hex=palette.warn, bold=True,
        )
        _textbox(
            slide, PAGE_MARGIN_L + Inches(0.35), y + Inches(0.48),
            CONTENT_W - Inches(0.7), card_h - Inches(0.6),
            f.get("reason", ""),
            font_name=palette.body_font, size_pt=12,
            colour_hex=palette.ink, line_spacing=1.35,
        )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Slide: Handover / closing
# ---------------------------------------------------------------------------

def build_closing(prs: Presentation, plan: dict, palette: Palette,
                  logo_path: Optional[Path], page_number: int, page_total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _page_background(slide, palette)
    proj = plan.get("project", {})

    # Accent kicker
    _rect(slide, PAGE_MARGIN_L, Inches(0.55), Inches(0.6), Inches(0.05),
          fill_hex=palette.accent, line_hex=None)

    _textbox(
        slide, PAGE_MARGIN_L, Inches(0.75),
        CONTENT_W, Inches(0.3),
        "HANDOVER",
        font_name=palette.body_font, size_pt=11,
        colour_hex=palette.ink_faint, bold=True,
    )
    _textbox(
        slide, PAGE_MARGIN_L, Inches(1.5),
        CONTENT_W, Inches(2),
        "Questions, or somewhere you'd like me to dig deeper?",
        font_name=palette.display_font, size_pt=42,
        colour_hex=palette.primary, bold=True,
        line_spacing=1.1,
    )

    _textbox(
        slide, PAGE_MARGIN_L, Inches(3.6),
        CONTENT_W * 0.8, Inches(1.5),
        "This deck pairs with a spreadsheet and an interactive HTML view "
        "covering the same plan at full resolution. Both ship with the "
        "package for your records.",
        font_name=palette.display_font, size_pt=18,
        colour_hex=palette.ink_soft, italic=True,
        line_spacing=1.4,
    )

    # Contact card bottom left
    card_y = Inches(5.3)
    card_w = Inches(4.2)
    card_h = Inches(1.25)
    _round_rect(slide, PAGE_MARGIN_L, card_y, card_w, card_h,
                fill_hex=palette.card, line_hex=palette.hairline,
                corner=0.06)
    _textbox(
        slide, PAGE_MARGIN_L + Inches(0.3), card_y + Inches(0.18),
        card_w - Inches(0.6), Inches(0.22),
        "PROJECT MANAGER",
        font_name=palette.body_font, size_pt=9,
        colour_hex=palette.ink_faint, bold=True,
    )
    _textbox(
        slide, PAGE_MARGIN_L + Inches(0.3), card_y + Inches(0.42),
        card_w - Inches(0.6), Inches(0.4),
        proj.get("pm", "—"),
        font_name=palette.display_font, size_pt=20,
        colour_hex=palette.primary, bold=True,
    )
    _textbox(
        slide, PAGE_MARGIN_L + Inches(0.3), card_y + Inches(0.85),
        card_w - Inches(0.6), Inches(0.28),
        f"{proj.get('client', '')} · {proj.get('name', '')}",
        font_name=palette.body_font, size_pt=10,
        colour_hex=palette.ink_soft, italic=True,
    )

    _footer(slide, palette=palette, project=proj, logo_path=logo_path,
            page_number=page_number, page_total=page_total)
    return slide


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------

def build_deck(plan: dict, plan_file: Path, out_path: Path) -> None:
    brand = plan.get("brand", {}) or {}
    palette = derive_palette(
        primary=brand.get("primary"),
        accent=brand.get("accent"),
        display_font=brand.get("display_font"),
        body_font=brand.get("body_font"),
    )
    logo_path = resolve_logo(plan.get("project", {}).get("logo_path"), plan_file)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    summary = _compute_summary(plan)
    has_flags = bool(plan.get("flags"))

    # Decide slide count up-front so footer page counters are correct
    slide_count = 7 + (1 if has_flags else 0)

    p = 1
    build_cover(prs, plan, palette, logo_path, p, slide_count); p += 1
    build_exec_snapshot(prs, plan, palette, logo_path, p, slide_count, summary); p += 1
    build_timeline_overview(prs, plan, palette, logo_path, p, slide_count, summary); p += 1
    build_detail_gantt(prs, plan, palette, logo_path, p, slide_count, summary); p += 1
    build_milestones(prs, plan, palette, logo_path, p, slide_count, summary); p += 1
    build_next_14_days(prs, plan, palette, logo_path, p, slide_count, summary); p += 1
    if has_flags:
        build_review_needed(prs, plan, palette, logo_path, p, slide_count); p += 1
    build_closing(prs, plan, palette, logo_path, p, slide_count)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    ap = argparse.ArgumentParser(
        description="Generate a premium-editorial client project-plan PPTX from a plan JSON."
    )
    ap.add_argument("--plan", required=True, help="Path to the plan JSON.")
    ap.add_argument("--out", required=True, help="Path to write the PPTX to.")
    args = ap.parse_args()

    plan_file = Path(args.plan).expanduser().resolve()
    out_file = Path(args.out).expanduser().resolve()
    with open(plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    build_deck(plan, plan_file, out_file)
    print(f"Wrote {out_file}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
